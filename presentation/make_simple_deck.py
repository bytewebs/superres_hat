#!/usr/bin/env python3
"""Generate SenHAT_Simple_Overview.pptx — broad, easy-to-understand deck."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

OUT = "SenHAT_Simple_Overview.pptx"
DIAGRAM = "senhat_two_phase.png"

C_BG = RGBColor(250, 251, 253)
C_TITLE = RGBColor(26, 61, 92)
C_TEXT = RGBColor(50, 50, 50)
C_P1 = RGBColor(46, 125, 90)
C_P2 = RGBColor(196, 92, 38)
C_ACCENT = RGBColor(0, 105, 148)


def bg(s):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG


def title(s, txt, y=Inches(0.4)):
    b = s.shapes.add_textbox(Inches(0.6), y, Inches(12.2), Inches(0.8))
    p = b.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = txt
    r.font.name = "Calibri"
    r.font.size = Pt(32)
    r.font.bold = True
    r.font.color.rgb = C_TITLE


def subtitle(s, txt, y=Inches(1.05), color=C_ACCENT):
    b = s.shapes.add_textbox(Inches(0.65), y, Inches(11.5), Inches(0.5))
    p = b.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = txt
    r.font.name = "Calibri"
    r.font.size = Pt(20)
    r.font.color.rgb = color


def bullets(s, items, y=Inches(1.35), sz=22):
    b = s.shapes.add_textbox(Inches(0.75), y, Inches(11.8), Inches(5.3))
    tf = b.text_frame
    tf.word_wrap = True
    for i, t in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = t
        p.font.name = "Calibri"
        p.font.size = Pt(sz)
        p.font.color.rgb = C_TEXT
        p.space_after = Pt(18)


def phase_box(s, left, top, w, h, phase, heading, items, color):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.color.rgb = C_TITLE
    sh.line.width = Pt(1.5)
    tf = sh.text_frame
    tf.clear()
    p0 = tf.paragraphs[0]
    p0.text = phase
    p0.font.bold = True
    p0.font.size = Pt(18)
    p0.font.color.rgb = C_TITLE
    p1 = tf.add_paragraph()
    p1.text = heading
    p1.font.size = Pt(14)
    p1.font.color.rgb = C_TEXT
    for it in items:
        p = tf.add_paragraph()
        p.text = "• " + it
        p.font.size = Pt(13)
        p.level = 0


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ---- Slide 1: Title ----
    s = prs.slides.add_slide(blank); bg(s)
    title(s, "SenHAT: Sentinel-2 Super-Resolution", Inches(1.2))
    subtitle(s, "A simple two-phase approach to accurate and sharp maps", Inches(2.1))
    b = s.shapes.add_textbox(Inches(0.65), Inches(5.7), Inches(8), Inches(0.5))
    b.text_frame.paragraphs[0].text = "Aayush Kharbanda  |  Meeting with Dr. Mishra"
    b.text_frame.paragraphs[0].font.size = Pt(18)

    # ---- Slide 2: Problem (simple) ----
    s = prs.slides.add_slide(blank); bg(s)
    title(s, "What problem are we solving?")
    bullets(s, [
        "Sentinel-2 gives useful data, but some bands are lower resolution (20 m) than others (10 m).",
        "We want one sharp 10 m map that keeps correct colours and vegetation signals.",
        "If the model invents fake edges or wrong band relationships, the map is not scientifically usable.",
    ], sz=21)

    # ---- Slide 3: SenHAT in plain language ----
    s = prs.slides.add_slide(blank); bg(s)
    title(s, "What is SenHAT?")
    bullets(s, [
        "SenHAT combines sharp 10 m geometry with richer 20 m spectral information.",
        "A transformer backbone looks at larger context (fields, rivers, roads) — not just tiny patches.",
        "Output: an 8-band super-resolved image that matches real Sentinel/VENµS structure.",
    ], sz=21)

    # ---- Slide 4: Why two phases ----
    s = prs.slides.add_slide(blank); bg(s)
    title(s, "Why train in two phases?")
    bullets(s, [
        "Phase 1 — Fidelity first: teach the model to get the right pixels, bands, and vegetation index.",
        "Phase 2 — Sharpness second: add perceptual sharpening only after the base map is trustworthy.",
        "This order reduces hallucination: we sharpen real structure, not invent fake detail.",
    ], sz=21)

    # ---- Slide 5: Phase 1 detail ----
    s = prs.slides.add_slide(blank); bg(s)
    title(s, "Phase 1: Fidelity Foundation")
    subtitle(s, "Current Narval run (no adversarial training yet)", y=Inches(1.0), color=C_P1)
    bullets(s, [
        "Losses: pixel (L1), content (VGG), artifact map (LDL), vegetation index (NDI).",
        "Discriminator is off — we are not chasing “pretty” textures yet.",
        "Success looks like: stable training, correct band mixing, boundaries that do not bleed.",
    ], y=Inches(1.45), sz=20)

    # ---- Slide 6: Phase 2 detail ----
    s = prs.slides.add_slide(blank); bg(s)
    title(s, "Phase 2: Controlled Sharpening")
    subtitle(s, "After fidelity weights are extracted", y=Inches(1.0), color=C_P2)
    bullets(s, [
        "Turn on adversarial + MS-SSIM to recover fine edges (roads, tree lines, field boundaries).",
        "Keep NDVI and spectral constraints active so sharpening does not break physics.",
        "Same SenHAT model — only the training objective changes.",
    ], y=Inches(1.45), sz=20)

    # ---- Slide 7: Two-phase diagram ----
    s = prs.slides.add_slide(blank); bg(s)
    title(s, "Two-Phase Training Pipeline")
    try:
        s.shapes.add_picture(DIAGRAM, Inches(0.5), Inches(1.05), width=Inches(12.3))
    except FileNotFoundError:
        sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2),
                                  Inches(12.3), Inches(5.5))
        sh.text_frame.paragraphs[0].text = "[Two-phase diagram]"

    # ---- Slide 8: Status & next steps ----
    s = prs.slides.add_slide(blank); bg(s)
    title(s, "Where we are now")
    bullets(s, [
        "Phase 1 fidelity training is running on Narval (A100, staged data on node SSD).",
        "Early signs: model learning structure; SSIM around 0.71 in initial checks.",
        "Next: pick best Phase 1 checkpoint → run Phase 2 sharpening → compare baselines for paper.",
    ], sz=20)

    prs.save(OUT)
    print(f"Created {OUT}")


if __name__ == "__main__":
    main()
