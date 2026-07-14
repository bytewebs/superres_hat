#!/usr/bin/env python3
"""Generate SenHAT_Final_Architecture.pptx"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

OUT = "SenHAT_Final_Architecture.pptx"

C_BG = RGBColor(248, 250, 253)
C_TITLE = RGBColor(22, 58, 92)
C_TEXT = RGBColor(45, 45, 45)
C_ACCENT = RGBColor(0, 110, 150)
C_PH = RGBColor(220, 228, 238)
C_PH_B = RGBColor(130, 150, 175)

W, H = Inches(13.333), Inches(7.5)


def bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_BG


def title(slide, txt, y=Inches(0.35)):
    b = slide.shapes.add_textbox(Inches(0.65), y, Inches(12.0), Inches(0.85))
    p = b.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = txt
    r.font.name = "Calibri"
    r.font.size = Pt(34)
    r.font.bold = True
    r.font.color.rgb = C_TITLE


def subtitle(slide, txt, y=Inches(1.15)):
    b = slide.shapes.add_textbox(Inches(0.7), y, Inches(11.5), Inches(0.55))
    p = b.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = txt
    r.font.name = "Calibri"
    r.font.size = Pt(22)
    r.font.color.rgb = C_ACCENT


def bullets(slide, items, y=Inches(1.4), sz=21):
    b = slide.shapes.add_textbox(Inches(0.85), y, Inches(11.6), Inches(5.2))
    tf = b.text_frame
    tf.word_wrap = True
    for i, t in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = t
        p.font.name = "Calibri"
        p.font.size = Pt(sz)
        p.font.color.rgb = C_TEXT
        p.space_after = Pt(16)


def placeholder(slide, l, t, w, h, label):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = C_PH
    s.line.color.rgb = C_PH_B
    s.line.width = Pt(2)
    p = s.text_frame.paragraphs[0]
    p.text = label
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Calibri"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = C_TITLE


def main():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    blank = prs.slide_layouts[6]

    # 1 Title
    s = prs.slides.add_slide(blank); bg(s)
    title(s, "SenHAT: Architecting Multispectral Super-Resolution for Sentinel-2", Inches(1.0))
    subtitle(s, "Fusing Spatial Geometry with Spectral Physics", Inches(2.2))
    b = s.shapes.add_textbox(Inches(0.7), Inches(5.8), Inches(8), Inches(0.5))
    b.text_frame.paragraphs[0].text = "Presenter: Aayush Kharbanda"
    b.text_frame.paragraphs[0].font.size = Pt(20)

    # 2 Problem
    s = prs.slides.add_slide(blank); bg(s)
    title(s, "The Hallucination Problem in Remote Sensing")
    bullets(s, [
        "Standard deep learning models (designed for natural images) rely on unconstrained generative priors, causing them to hallucinate geometry.",
        "In satellite mapping, hallucinating a building or shifting a riverbank destroys the scientific validity of the dataset.",
        "Goal: Engineer an architecture that understands macro-geography while strictly obeying multispectral physics.",
    ])

    # 3 Architecture
    s = prs.slides.add_slide(blank); bg(s)
    title(s, "SenHAT Architecture Breakdown")
    b = s.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(5.0), Inches(5.0))
    tf = b.text_frame; tf.word_wrap = True
    for i, t in enumerate([
        "Fuses high-frequency 10m spatial geometry with 20m spectral data prior to deep extraction.",
        "Replaces localized CNNs with an Overlapping Cross-Attention Transformer to capture global geographical context.",
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = t; p.font.size = Pt(20); p.space_after = Pt(18)
    placeholder(s, Inches(6.0), Inches(1.15), Inches(6.8), Inches(5.8),
                "[Drop in LaTeX Architecture Diagram Here]")

    # 4 Guardrails
    s = prs.slides.add_slide(blank); bg(s)
    title(s, "Constraining the Generation")
    bullets(s, [
        "Unlike standard perceptual models, our network is mathematically penalized for violating terrestrial physics.",
        "NDVI Loss ensures the biological reality of the vegetation index is preserved.",
        "Spectral Correlation (SCC) Loss locks the relationships between the Red Edge and SWIR bands, preventing color distortion.",
    ])

    # 5 Adversarial sharpening
    s = prs.slides.add_slide(blank); bg(s)
    title(s, "Adversarial Texture Optimization (Without Hallucination)")
    bullets(s, [
        "Traditional Mean Squared Error (MSE/L1) training inherently smooths out high-frequency spatial edges (like roads and tree lines).",
        "To recover photorealistic sharpness, we utilize a secondary adversarial network to optimize perceptual textures.",
        "Because this adversarial tuning is constrained by our NDVI and SCC guardrails, the model sharpens actual geometry instead of inventing fake textures.",
    ])

    # 6 Results
    s = prs.slides.add_slide(blank); bg(s)
    title(s, "Early Spatial vs. Spectral Fidelity")
    placeholder(s, Inches(0.65), Inches(1.1), Inches(12.0), Inches(5.4),
                "[Drop in Narval Training Grid PNG Here]")
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(6.6), Inches(12.0), Inches(0.55))
    b.fill.solid(); b.fill.fore_color.rgb = RGBColor(230, 236, 245)
    b.line.color.rgb = C_PH_B
    p = b.text_frame.paragraphs[0]
    p.text = ("The architecture is successfully utilizing 10m spatial conditioning to prevent "
              "spectral bleeding across agricultural boundaries. Early SSIM hitting 0.71.")
    p.font.size = Pt(15)

    # 7 Next steps
    s = prs.slides.add_slide(blank); bg(s)
    title(s, "Timeline to Publication")
    bullets(s, [
        "Extract optimal weights from the current fidelity run.",
        "Execute the adversarial perceptual optimization phase.",
        "Generate precision-recall curves and baseline comparisons for ISPRS journal submission.",
    ])

    prs.save(OUT)
    print(f"Created {OUT}")


if __name__ == "__main__":
    main()
